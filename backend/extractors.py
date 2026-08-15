"""Turn whatever the user supplies into plain text.

Supported: .pdf .docx .pptx .txt .md .csv, article URLs, PDFs linked by URL,
and YouTube videos with captions.

Every failure raises ExtractionError carrying a message that is safe and
useful to show a student.
"""
from __future__ import annotations

import io
import logging
import re
from urllib.parse import parse_qs, urlparse

import httpx

from config import get_settings

log = logging.getLogger("dvq.extract")
settings = get_settings()

# A bare user-agent is not enough: WordPress security plugins and WAFs check
# for the full set of headers a real browser sends.
BROWSER_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
        "(KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36"
    ),
    "Accept": ("text/html,application/xhtml+xml,application/xml;q=0.9,"
               "image/avif,image/webp,*/*;q=0.8"),
    "Accept-Language": "en-IN,en-GB;q=0.9,en;q=0.8,hi;q=0.7",
    "Upgrade-Insecure-Requests": "1",
    "Sec-Fetch-Dest": "document",
    "Sec-Fetch-Mode": "navigate",
    "Sec-Fetch-Site": "none",
    "Sec-Fetch-User": "?1",
    "Cache-Control": "max-age=0",
    "Connection": "keep-alive",
}


class ExtractionError(Exception):
    """We could not get usable text out of the source."""


# ---------------------------------------------------------------------------
# Files
# ---------------------------------------------------------------------------
def from_pdf(data: bytes) -> str:
    import fitz  # PyMuPDF

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError("That PDF could not be opened. It may be corrupt.") from exc

    try:
        parts = [page.get_text("text") for page in doc]
        pages = doc.page_count
    finally:
        doc.close()

    text = "\n".join(parts)
    if len(text.strip()) < settings.MIN_SOURCE_CHARS:
        raise ExtractionError(
            f"This PDF has almost no selectable text across its {pages} page(s) — "
            f"it is probably a scan. Run it through OCR first, or paste the text."
        )
    return text


def from_docx(data: bytes) -> str:
    from docx import Document

    try:
        doc = Document(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError("That Word file could not be opened.") from exc

    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def from_pptx(data: bytes) -> str:
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError("That PowerPoint file could not be opened.") from exc

    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
        notes = getattr(slide, "notes_slide", None)
        if notes is not None and notes.notes_text_frame is not None:
            note = notes.notes_text_frame.text.strip()
            if note:
                parts.append(f"[Notes] {note}")
    return "\n".join(parts)


def from_plaintext(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "cp1252", "latin-1"):
        try:
            return data.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="ignore")


EXT_HANDLERS = {
    ".pdf": from_pdf,
    ".docx": from_docx,
    ".pptx": from_pptx,
}

PLAIN_EXTS = (".txt", ".md", ".markdown", ".csv", ".tsv", ".rtf", ".log")


def sniff_kind(filename: str, data: bytes) -> str:
    """Trust the magic bytes over the extension."""
    if data[:5] == b"%PDF-":
        return "pdf"
    if data[:4] == b"PK\x03\x04":                       # any OOXML zip
        head = data[:4000]
        if b"word/" in head:
            return "docx"
        if b"ppt/" in head:
            return "pptx"
    name = (filename or "").lower()
    for ext in EXT_HANDLERS:
        if name.endswith(ext):
            return ext.lstrip(".")
    if name.endswith(PLAIN_EXTS):
        return "text"
    return ""


def from_file(filename: str, data: bytes) -> str:
    if not data:
        raise ExtractionError("That file is empty.")

    kind = sniff_kind(filename, data)
    if kind == "pdf":
        return from_pdf(data)
    if kind == "docx":
        return from_docx(data)
    if kind == "pptx":
        return from_pptx(data)
    if kind == "text":
        return from_plaintext(data)

    # Last resort: if it decodes to sensible text, take it.
    try:
        guess = data.decode("utf-8")
        if sum(c.isprintable() or c.isspace() for c in guess[:2000]) > len(guess[:2000]) * 0.9:
            return guess
    except UnicodeDecodeError:
        pass

    raise ExtractionError(
        "Unsupported file type. Upload a PDF, DOCX, PPTX, TXT or MD file."
    )


# ---------------------------------------------------------------------------
# Web pages
# ---------------------------------------------------------------------------
BLOCKED_STATUS = {401: "requires a login", 403: "blocked the request",
                  406: "refused the request", 429: "asked us to slow down",
                  451: "is not available for legal reasons"}

STRIP_TAGS = re.compile(
    r"<(script|style|nav|header|footer|form|noscript|svg|aside)\b[^>]*>.*?</\1>",
    re.I | re.S,
)
TAGS = re.compile(r"<[^>]+>")
WS = re.compile(r"[ \t\u00a0]+")


def _html_to_text(html: str) -> str:
    """Crude fallback when the article extractor finds nothing."""
    text = STRIP_TAGS.sub(" ", html)
    text = re.sub(r"<br\s*/?>|</p>|</div>|</li>|</h[1-6]>", "\n", text, flags=re.I)
    text = TAGS.sub(" ", text)
    for entity, char in (("&nbsp;", " "), ("&amp;", "&"), ("&lt;", "<"),
                         ("&gt;", ">"), ("&quot;", '"'), ("&#39;", "'")):
        text = text.replace(entity, char)
    lines = [WS.sub(" ", ln).strip() for ln in text.split("\n")]
    return "\n".join(ln for ln in lines if len(ln) > 2)


def normalise_url(url: str) -> str:
    url = (url or "").strip()
    if not url:
        raise ExtractionError("Enter a link first.")
    # A scheme is anything before a colon that is not a port number.
    scheme = re.match(r"^([a-zA-Z][a-zA-Z0-9+.\-]*):(?!\d)", url)
    if scheme and scheme.group(1).lower() not in ("http", "https"):
        raise ExtractionError("Only http and https links are supported.")
    if "://" not in url:
        url = "https://" + url.lstrip("/")
    parsed = urlparse(url)
    if parsed.scheme not in ("http", "https"):
        raise ExtractionError("Only http and https links are supported.")
    if not parsed.netloc:
        raise ExtractionError("That does not look like a valid link.")
    return url


def fetch_url(url: str) -> httpx.Response:
    try:
        with httpx.Client(timeout=settings.FETCH_TIMEOUT, follow_redirects=True,
                          headers=BROWSER_HEADERS) as client:
            return client.get(url)
    except httpx.TimeoutException as exc:
        raise ExtractionError(
            f"That page took longer than {settings.FETCH_TIMEOUT}s to respond."
        ) from exc
    except httpx.TooManyRedirects as exc:
        raise ExtractionError("That link redirects in a loop.") from exc
    except httpx.ConnectError as exc:
        raise ExtractionError(
            "That site could not be reached. Check the address is spelled correctly."
        ) from exc
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError("That page could not be loaded.") from exc


def from_url(url: str) -> str:
    url = normalise_url(url)
    resp = fetch_url(url)

    if resp.status_code in BLOCKED_STATUS:
        raise ExtractionError(
            f"That site {BLOCKED_STATUS[resp.status_code]} (HTTP {resp.status_code}). "
            f"Its firewall blocks automated readers — paste the text instead."
        )
    if resp.status_code == 404:
        raise ExtractionError("That page does not exist (HTTP 404).")
    if resp.status_code >= 400:
        raise ExtractionError(f"That page returned HTTP {resp.status_code}.")

    content_type = resp.headers.get("content-type", "").lower()

    # A link straight to a PDF is common for syllabus and PYQ papers.
    if "application/pdf" in content_type or resp.content[:5] == b"%PDF-":
        return from_pdf(resp.content)

    if content_type.startswith("text/plain"):
        return resp.text

    if content_type and not any(t in content_type for t in ("html", "xml", "text")):
        raise ExtractionError(
            f"That link points at a {content_type.split(';')[0]} file, "
            f"which cannot be read as text."
        )

    html = resp.text
    text = ""

    try:
        import trafilatura
        text = trafilatura.extract(
            html, include_comments=False, include_tables=True,
            favor_recall=True, url=url,
        ) or ""
    except Exception as exc:  # noqa: BLE001
        log.warning("trafilatura failed on %s: %s", url, exc)

    if len(text.strip()) < settings.MIN_SOURCE_CHARS:
        fallback = _html_to_text(html)
        if len(fallback.strip()) > len(text.strip()):
            log.info("used html fallback for %s", url)
            text = fallback

    if len(text.strip()) < settings.MIN_SOURCE_CHARS:
        raise ExtractionError(
            "Not enough readable text on that page. It may be mostly images, "
            "or built entirely in JavaScript. Paste the content instead."
        )
    return text


# ---------------------------------------------------------------------------
# YouTube
# ---------------------------------------------------------------------------
YT_ID = re.compile(r"^[A-Za-z0-9_-]{11}$")


def youtube_id(url: str) -> str:
    url = (url or "").strip()
    parsed = urlparse(url if "//" in url else "https://" + url)
    host = parsed.netloc.lower()

    if host.endswith("youtu.be"):
        candidate = parsed.path.lstrip("/").split("/")[0]
    elif "youtube" in host:
        if parsed.path.startswith(("/shorts/", "/embed/", "/live/", "/v/")):
            parts = [p for p in parsed.path.split("/") if p]
            candidate = parts[1] if len(parts) > 1 else ""
        else:
            candidate = parse_qs(parsed.query).get("v", [""])[0]
    else:
        candidate = url

    candidate = candidate.split("?")[0].split("&")[0].strip()
    if not YT_ID.match(candidate):
        raise ExtractionError(
            "That does not look like a YouTube video link. Paste the full URL "
            "from the address bar."
        )
    return candidate


def _snippet_text(chunk) -> str:
    """0.6.x yields dicts, 1.x yields snippet objects."""
    if isinstance(chunk, dict):
        return chunk.get("text", "")
    return getattr(chunk, "text", "")


def from_youtube(url: str, languages: tuple[str, ...] = ("en", "en-IN", "hi")) -> str:
    vid = youtube_id(url)          # validate before touching the library

    try:
        from youtube_transcript_api import YouTubeTranscriptApi
    except ImportError as exc:
        raise ExtractionError(
            "YouTube support is not installed on this server. "
            "Run: pip install youtube-transcript-api"
        ) from exc

    wanted = list(languages)
    errors: list[str] = []

    def join(chunks) -> str:
        return " ".join(_snippet_text(c) for c in chunks)

    # --- 1.x: instance methods -------------------------------------------
    if hasattr(YouTubeTranscriptApi, "fetch"):
        api = YouTubeTranscriptApi()
        try:
            return join(api.fetch(vid, languages=wanted))
        except Exception as exc:  # noqa: BLE001
            errors.append(f"{type(exc).__name__}: {exc}")
        try:
            for transcript in api.list(vid):
                try:
                    return join(transcript.fetch())
                except Exception:  # noqa: BLE001
                    continue
        except Exception as exc:  # noqa: BLE001
            errors.append(f"{type(exc).__name__}: {exc}")

    # --- 0.6.x: class methods --------------------------------------------
    if hasattr(YouTubeTranscriptApi, "get_transcript"):
        try:
            return join(YouTubeTranscriptApi.get_transcript(vid, languages=wanted))
        except Exception as exc:  # noqa: BLE001
            errors.append(f"{type(exc).__name__}: {exc}")
        try:
            for transcript in YouTubeTranscriptApi.list_transcripts(vid):
                try:
                    return join(transcript.fetch())
                except Exception:  # noqa: BLE001
                    continue
        except Exception as exc:  # noqa: BLE001
            errors.append(f"{type(exc).__name__}: {exc}")

    detail = errors[0] if errors else "no transcript tracks were returned"
    log.warning("youtube %s failed: %s", vid, detail)

    if "IpBlocked" in detail or "blocked" in detail.lower():
        raise ExtractionError(
            "YouTube is blocking transcript requests from this server. This is "
            "common on cloud hosting. Paste the transcript text instead."
        )
    raise ExtractionError(
        "No captions could be read for this video. Try one with subtitles turned "
        f"on, or paste the text instead. ({detail[:120]})"
    )


# ---------------------------------------------------------------------------
# Normalise
# ---------------------------------------------------------------------------
def condense(text: str, limit: int | None = None) -> str:
    """Collapse whitespace, then keep head + middle + tail if over the limit.

    Sampling three regions beats a plain head-truncate: questions then cover
    the whole document instead of only chapter one.
    """
    limit = limit or settings.max_source_chars
    text = WS.sub(" ", text or "")
    text = re.sub(r"\n{3,}", "\n\n", text).strip()

    if len(text) <= limit:
        return text

    slice_len = limit // 3
    head = text[:slice_len]
    mid_start = (len(text) - slice_len) // 2
    mid = text[mid_start:mid_start + slice_len]
    tail = text[-slice_len:]
    return f"{head}\n\n[...]\n\n{mid}\n\n[...]\n\n{tail}"


def split_sections(text: str, count: int) -> list[str]:
    """Split into `count` roughly equal parts on paragraph boundaries.

    Used for batched generation so each batch of questions comes from a
    different part of the document.
    """
    text = text.strip()
    if count <= 1 or len(text) < 800:
        return [text]

    paragraphs = [p for p in re.split(r"\n\s*\n", text) if p.strip()]
    if len(paragraphs) < count:
        size = max(1, len(text) // count)
        return [text[i:i + size] for i in range(0, len(text), size)][:count]

    per = len(paragraphs) / count
    out: list[str] = []
    for i in range(count):
        chunk = paragraphs[int(i * per):int((i + 1) * per)]
        if chunk:
            out.append("\n\n".join(chunk))
    return out or [text]
